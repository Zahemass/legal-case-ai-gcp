# services/extraction-service/src/extractor.py
import os
import time
import logging
import tempfile
from typing import Dict, Any, Optional
import chardet
import magic

# PDF processing
import PyPDF2
from pdf2image import convert_from_path

# Office document processing
from docx import Document
from pptx import Presentation
import openpyxl
import xlrd

# OCR processing
try:
    import pytesseract
    from PIL import Image
    OCR_AVAILABLE = True
except ImportError:
    OCR_AVAILABLE = False
    logging.warning("OCR libraries not available - image processing will be limited")

logger = logging.getLogger(__name__)

class DocumentExtractor:
    """Document text extraction service supporting multiple file formats"""
    
    def __init__(self):
        self.supported_types = {
            'application/pdf': self._extract_pdf,
            'application/msword': self._extract_doc_legacy,
            'application/vnd.openxmlformats-officedocument.wordprocessingml.document': self._extract_docx,
            'application/vnd.ms-powerpoint': self._extract_ppt_legacy,
            'application/vnd.openxmlformats-officedocument.presentationml.presentation': self._extract_pptx,
            'application/vnd.ms-excel': self._extract_xls_legacy,
            'application/vnd.openxmlformats-officedocument.spreadsheetml.sheet': self._extract_xlsx,
            'text/plain': self._extract_text,
            'text/csv': self._extract_text,
            'application/rtf': self._extract_rtf,
            'image/jpeg': self._extract_image_ocr,
            'image/png': self._extract_image_ocr,
            'image/gif': self._extract_image_ocr,
            'image/webp': self._extract_image_ocr,
            'image/tiff': self._extract_image_ocr
        }
        
        logger.info(f"✅ DocumentExtractor initialized with {len(self.supported_types)} supported types")

    def extract_text(self, file_path: str, content_type: str, filename: str = None) -> Dict[str, Any]:
        """Extract text from various document types"""
        start_time = time.time()
        
        try:
            logger.info(f"🔍 Starting extraction: {filename} ({content_type})")
            
            if not os.path.exists(file_path):
                raise FileNotFoundError(f"File not found: {file_path}")
            
            file_size = os.path.getsize(file_path)
            if file_size == 0:
                raise ValueError("File is empty")
            
            # Check if content type is supported
            if content_type not in self.supported_types:
                raise ValueError(f"Unsupported content type: {content_type}")
            
            # Extract text using appropriate method
            extractor_method = self.supported_types[content_type]
            result = extractor_method(file_path, filename)
            
            # Add common metadata
            processing_time = time.time() - start_time
            result.update({
                'content_type': content_type,
                'file_size': file_size,
                'processing_time': processing_time,
                'method': extractor_method.__name__,
                'extracted_at': time.time()
            })
            
            if not result.get('text'):
                result['text'] = ''
            
            text_length = len(result['text'])
            logger.info(f"✅ Extraction completed: {text_length} characters in {processing_time:.2f}s")
            
            return result
            
        except Exception as e:
            processing_time = time.time() - start_time
            logger.error(f"❌ Extraction failed for {filename}: {str(e)}")
            
            return {
                'text': '',
                'page_count': 0,
                'title': filename or 'Unknown',
                'language': 'unknown',
                'confidence': 0.0,
                'error': str(e),
                'processing_time': processing_time,
                'method': 'error'
            }

    def _extract_pdf(self, file_path: str, filename: str = None) -> Dict[str, Any]:
        """Extract text from PDF files"""
        try:
            text = ""
            page_count = 0
            title = ""
            
            with open(file_path, 'rb') as file:
                pdf_reader = PyPDF2.PdfReader(file)
                page_count = len(pdf_reader.pages)
                
                if pdf_reader.metadata:
                    title = pdf_reader.metadata.get('/Title', '') or filename or ''
                
                for page in pdf_reader.pages:
                    page_text = page.extract_text()
                    if page_text:
                        text += page_text + "\n\n"
            
            if not title and text:
                first_line = text.split('\n')[0].strip()
                title = first_line[:100] if len(first_line) > 0 else filename or 'PDF Document'
            
            return {
                'text': text.strip(),
                'page_count': page_count,
                'title': title or filename or 'PDF Document',
                'language': self._detect_language(text),
                'confidence': 0.9 if text.strip() else 0.1,
                'metadata': {
                    'pages': page_count,
                    'extraction_method': 'PyPDF2'
                }
            }
            
        except Exception as e:
            logger.error(f"❌ PDF extraction error: {e}")
            raise e

    def _extract_docx(self, file_path: str, filename: str = None) -> Dict[str, Any]:
        """Extract text from DOCX files"""
        try:
            doc = Document(file_path)
            text = ""
            
            for paragraph in doc.paragraphs:
                if paragraph.text.strip():
                    text += paragraph.text + "\n"
            
            for table in doc.tables:
                for row in table.rows:
                    row_text = []
                    for cell in row.cells:
                        if cell.text.strip():
                            row_text.append(cell.text.strip())
                    if row_text:
                        text += " | ".join(row_text) + "\n"
            
            title = filename or 'Word Document'
            if doc.paragraphs and doc.paragraphs[0].text.strip():
                first_para = doc.paragraphs[0].text.strip()
                if len(first_para) < 100:
                    title = first_para
            
            return {
                'text': text.strip(),
                'page_count': 1,
                'title': title,
                'language': self._detect_language(text),
                'confidence': 0.95,
                'metadata': {
                    'extraction_method': 'python-docx'
                }
            }
            
        except Exception as e:
            logger.error(f"❌ DOCX extraction error: {e}")
            raise e

    def _extract_xlsx(self, file_path: str, filename: str = None) -> Dict[str, Any]:
        """Extract text from XLSX files"""
        try:
            workbook = openpyxl.load_workbook(file_path, data_only=True)
            text = ""
            
            for sheet_name in workbook.sheetnames:
                sheet = workbook[sheet_name]
                text += f"--- Sheet: {sheet_name} ---\n"
                
                for row in sheet.iter_rows(values_only=True):
                    row_text = []
                    for cell_value in row:
                        if cell_value is not None:
                            row_text.append(str(cell_value))
                    
                    if row_text:
                        text += " | ".join(row_text) + "\n"
                
                text += "\n"
            
            return {
                'text': text.strip(),
                'page_count': len(workbook.worksheets),
                'title': filename or 'Excel Spreadsheet',
                'language': self._detect_language(text),
                'confidence': 0.9,
                'metadata': {
                    'sheets': len(workbook.worksheets),
                    'extraction_method': 'openpyxl'
                }
            }
            
        except Exception as e:
            logger.error(f"❌ XLSX extraction error: {e}")
            raise e

    def _extract_text(self, file_path: str, filename: str = None) -> Dict[str, Any]:
        """Extract text from plain text files"""
        try:
            with open(file_path, 'rb') as file:
                raw_data = file.read()
                encoding_result = chardet.detect(raw_data)
                encoding = encoding_result['encoding'] or 'utf-8'
            
            with open(file_path, 'r', encoding=encoding, errors='ignore') as file:
                text = file.read()
            
            return {
                'text': text,
                'page_count': max(1, text.count('\n') // 50),
                'title': filename or 'Text Document',
                'language': self._detect_language(text),
                'confidence': 0.95,
                'metadata': {
                    'encoding': encoding,
                    'extraction_method': 'text_file'
                }
            }
            
        except Exception as e:
            logger.error(f"❌ Text extraction error: {e}")
            raise e

    def _extract_pptx(self, file_path: str, filename: str = None) -> Dict[str, Any]:
        """Extract text from PPTX files"""
        try:
            prs = Presentation(file_path)
            text = ""
            
            for slide_num, slide in enumerate(prs.slides, 1):
                slide_text = f"--- Slide {slide_num} ---\n"
                
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        slide_text += shape.text + "\n"
                
                if slide_text.strip() != f"--- Slide {slide_num} ---":
                    text += slide_text + "\n"
            
            return {
                'text': text.strip(),
                'page_count': len(prs.slides),
                'title': filename or 'PowerPoint Presentation',
                'language': self._detect_language(text),
                'confidence': 0.9,
                'metadata': {
                    'slides': len(prs.slides),
                    'extraction_method': 'python-pptx'
                }
            }
            
        except Exception as e:
            logger.error(f"❌ PPTX extraction error: {e}")
            raise e

    def _extract_image_ocr(self, file_path: str, filename: str = None) -> Dict[str, Any]:
        """Extract text from images using OCR"""
        if not OCR_AVAILABLE:
            raise Exception("OCR libraries not available")
        
        try:
            image = Image.open(file_path)
            if image.mode != 'RGB':
                image = image.convert('RGB')
            
            text = pytesseract.image_to_string(image, lang='eng')
            
            return {
                'text': text.strip(),
                'page_count': 1,
                'title': filename or 'Image Document',
                'language': self._detect_language(text),
                'confidence': 0.6,
                'metadata': {
                    'ocr_engine': 'tesseract',
                    'extraction_method': 'OCR'
                }
            }
            
        except Exception as e:
            logger.error(f"❌ Image OCR extraction error: {e}")
            return {
                'text': '',
                'page_count': 1,
                'title': filename or 'Image Document',
                'language': 'unknown',
                'confidence': 0.0,
                'error': str(e)
            }

    def _extract_rtf(self, file_path: str, filename: str = None) -> Dict[str, Any]:
        """Extract text from RTF files"""
        try:
            with open(file_path, 'r', encoding='utf-8', errors='ignore') as file:
                content = file.read()
            
            # Basic RTF text extraction
            text = ""
            i = 0
            
            while i < len(content):
                char = content[i]
                
                if char == '\\':
                    # Skip RTF command
                    while i < len(content) and content[i] not in [' ', '\n', '}']:
                        i += 1
                elif char not in ['{', '}']:
                    text += char
                
                i += 1
            
            text = ' '.join(text.split())
            
            return {
                'text': text,
                'page_count': max(1, len(text) // 2000),
                'title': filename or 'RTF Document',
                'language': self._detect_language(text),
                'confidence': 0.7,
                'metadata': {
                    'extraction_method': 'basic_rtf_parser'
                }
            }
            
        except Exception as e:
            logger.error(f"❌ RTF extraction error: {e}")
            raise e

    def _extract_doc_legacy(self, file_path: str, filename: str = None) -> Dict[str, Any]:
        """Extract text from legacy DOC files"""
        raise Exception("Legacy DOC format not supported. Please convert to DOCX.")

    def _extract_ppt_legacy(self, file_path: str, filename: str = None) -> Dict[str, Any]:
        """Extract text from legacy PPT files"""
        raise Exception("Legacy PPT format not supported. Please convert to PPTX.")

    def _extract_xls_legacy(self, file_path: str, filename: str = None) -> Dict[str, Any]:
        """Extract text from legacy XLS files"""
        try:
            workbook = xlrd.open_workbook(file_path)
            text = ""
            
            for sheet_index in range(workbook.nsheets):
                sheet = workbook.sheet_by_index(sheet_index)
                text += f"--- Sheet: {sheet.name} ---\n"
                
                for row_num in range(sheet.nrows):
                    row_text = []
                    for col_num in range(sheet.ncols):
                        cell_value = sheet.cell_value(row_num, col_num)
                        if cell_value:
                            row_text.append(str(cell_value))
                    
                    if row_text:
                        text += " | ".join(row_text) + "\n"
            
            return {
                'text': text.strip(),
                'page_count': workbook.nsheets,
                'title': filename or 'Excel Spreadsheet (Legacy)',
                'language': self._detect_language(text),
                'confidence': 0.85,
                'metadata': {
                    'sheets': workbook.nsheets,
                    'extraction_method': 'xlrd'
                }
            }
            
        except Exception as e:
            logger.error(f"❌ XLS extraction error: {e}")
            raise e

    def _detect_language(self, text: str) -> str:
        """Simple language detection"""
        if not text or len(text.strip()) < 50:
            return 'unknown'
        
        text_lower = text.lower()
        english_words = ['the', 'and', 'or', 'but', 'in', 'on', 'at', 'to', 'for', 'of', 'with', 'by']
        english_count = sum(1 for word in english_words if f' {word} ' in text_lower)
        
        if english_count >= 3:
            return 'en'
        
        return 'unknown'

    def is_supported(self, content_type: str) -> bool:
        """Check if content type is supported"""
        return content_type in self.supported_types